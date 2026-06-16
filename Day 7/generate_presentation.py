"""
Day 7 - Generate 12-Slide Presentation Deck (PPTX)
Creates a professional PowerPoint presentation for the Bluestock MF Capstone.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.join(SCRIPT_DIR, '..')
OUTPUT_PATH = os.path.join(SCRIPT_DIR, 'Bluestock_MF_Presentation.pptx')
SCREENSHOT_DIR = os.path.join(PROJECT_ROOT, 'Day 5', 'Dashboard_img')

# Theme colors
DARK_BLUE = RGBColor(0, 51, 102)
ACCENT_BLUE = RGBColor(0, 102, 204)
LIGHT_BLUE = RGBColor(200, 220, 240)
LIGHT_GREY = RGBColor(220, 220, 220)
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(40, 40, 40)
GREY_TEXT = RGBColor(100, 100, 100)
BG_LIGHT = RGBColor(245, 248, 252)


def add_bg_shape(slide, prs, color=None):
    """Add a colored background rectangle to the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color or DARK_BLUE
    shape.line.fill.background()


def add_accent_bar(slide, prs):
    """Add a thick blue accent bar at the top."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()


def add_side_bar(slide):
    """Add a dark blue side accent bar on the left."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0.06),
        Inches(0.12), Inches(7.44)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()


def add_title_textbox(slide, text, left, top, width, height,
                      font_size=28, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT):
    """Add a styled title text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_body_textbox(slide, text, left, top, width, height,
                     font_size=14, color=DARK_TEXT, line_spacing=20):
    """Add a body text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.line_spacing = Pt(line_spacing)
    return tf


def add_content_slide(slide, prs, title, bullets, extra_paragraphs=None):
    """Add a fully packed content slide with title, bullets, and optional extra text."""
    add_accent_bar(slide, prs)
    add_side_bar(slide)

    # Title
    add_title_textbox(slide, title, Inches(0.5), Inches(0.2), Inches(9), Inches(0.7),
                      font_size=26, color=DARK_BLUE)

    # Divider line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(9), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Bullet content
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (bold_text, normal_text) in enumerate(bullets):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.space_after = Pt(6)
        p.space_before = Pt(4)
        p.line_spacing = Pt(20)

        if bold_text:
            run_b = p.add_run()
            run_b.text = bold_text
            run_b.font.bold = True
            run_b.font.size = Pt(13)
            run_b.font.color.rgb = DARK_BLUE

        run_n = p.add_run()
        run_n.text = normal_text
        run_n.font.size = Pt(13)
        run_n.font.color.rgb = DARK_TEXT

    # Extra paragraphs at the bottom if provided
    if extra_paragraphs:
        for para_text in extra_paragraphs:
            p = tf.add_paragraph()
            p.space_before = Pt(6)
            p.line_spacing = Pt(20)
            run = p.add_run()
            run.text = para_text
            run.font.size = Pt(12)
            run.font.color.rgb = GREY_TEXT
            run.font.italic = True

    # Footer bar
    footer_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.1),
        prs.slide_width, Inches(0.4)
    )
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = DARK_BLUE
    footer_bar.line.fill.background()

    ftBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.12), Inches(9), Inches(0.35))
    ft_tf = ftBox.text_frame
    ft_p = ft_tf.paragraphs[0]
    ft_p.text = "Bluestock MF Analytics | Damodara P | June 2026"
    ft_p.font.size = Pt(9)
    ft_p.font.color.rgb = LIGHT_GREY
    ft_p.alignment = PP_ALIGN.CENTER


def add_screenshot_slide(slide, prs, title, img_filename, caption):
    """Add a slide with a dashboard screenshot."""
    add_accent_bar(slide, prs)
    add_side_bar(slide)
    add_title_textbox(slide, title, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6),
                      font_size=22, color=DARK_BLUE)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), Inches(9), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    img_path = os.path.join(SCREENSHOT_DIR, img_filename)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.0), Inches(9), Inches(5.2))
    else:
        add_body_textbox(slide, f'[Screenshot not found: {img_filename}]',
                         Inches(2), Inches(3), Inches(6), Inches(1),
                         font_size=16, color=GREY_TEXT)

    add_body_textbox(slide, caption, Inches(0.5), Inches(6.3), Inches(9), Inches(0.5),
                     font_size=11, color=GREY_TEXT)

    # Footer bar
    footer_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.1),
        prs.slide_width, Inches(0.4)
    )
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = DARK_BLUE
    footer_bar.line.fill.background()

    ftBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.12), Inches(9), Inches(0.35))
    ft_tf = ftBox.text_frame
    ft_p = ft_tf.paragraphs[0]
    ft_p.text = "Bluestock MF Analytics | Damodara P | June 2026"
    ft_p.font.size = Pt(9)
    ft_p.font.color.rgb = LIGHT_GREY
    ft_p.alignment = PP_ALIGN.CENTER


# ==================== BUILD PRESENTATION ====================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ====== SLIDE 1: TITLE ======
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg_shape(slide1, prs)

# Accent line
bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(1.2), Inches(5), Inches(0.04))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

add_title_textbox(slide1, 'BLUESTOCK MUTUAL FUND', Inches(0.5), Inches(1.5),
                  Inches(9), Inches(0.8), font_size=38, color=WHITE, alignment=PP_ALIGN.CENTER)
add_title_textbox(slide1, 'ANALYTICS CAPSTONE', Inches(0.5), Inches(2.3),
                  Inches(9), Inches(0.8), font_size=38, color=WHITE, alignment=PP_ALIGN.CENTER)

bar2 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(3.3), Inches(4), Inches(0.04))
bar2.fill.solid()
bar2.fill.fore_color.rgb = ACCENT_BLUE
bar2.line.fill.background()

add_title_textbox(slide1, 'End-to-End Data Engineering, Dashboarding & Risk Analysis',
                  Inches(0.5), Inches(3.6), Inches(9), Inches(0.6),
                  font_size=16, bold=False, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

add_title_textbox(slide1, 'Submitted By: Damodara P',
                  Inches(0.5), Inches(4.8), Inches(9), Inches(0.5),
                  font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_title_textbox(slide1, 'Data Analyst Intern  |  Bluestock Fintech',
                  Inches(0.5), Inches(5.3), Inches(9), Inches(0.5),
                  font_size=14, bold=False, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)
add_title_textbox(slide1, 'pdamodaran2000@gmail.com  |  June 2026',
                  Inches(0.5), Inches(5.8), Inches(9), Inches(0.5),
                  font_size=13, bold=False, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

# ====== SLIDE 2: PROBLEM & OBJECTIVE ======
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide2, prs, 'Problem Statement & Objectives', [
    ('The Challenge: ', 'The Indian Mutual Fund industry manages Rs.81 Lakh Crore across 1,900+ schemes. '
     'Retail investors lack accessible, data-driven tools to compare risk-adjusted fund performance '
     'and make informed investment decisions.'),
    ('Objective 1 - Data Pipeline: ', 'Build an automated ETL pipeline to ingest, clean, validate, '
     'and store 17 raw datasets (127,000+ rows) into a structured Star Schema SQLite database.'),
    ('Objective 2 - Performance Analytics: ', 'Calculate CAGR (1yr/3yr/5yr), Sharpe Ratio, Sortino Ratio, '
     'Alpha, Beta, Maximum Drawdown, VaR, and CVaR for 40 mutual fund schemes.'),
    ('Objective 3 - Interactive Dashboard: ', 'Develop a 4-page interactive Power BI dashboard with '
     'KPI cards, scatter plots, drill-through navigation, slicers, and tooltips.'),
    ('Objective 4 - Recommendation Engine: ', 'Build a Python-based algorithmic fund recommender '
     'that maps investor risk profiles (Low/Moderate/High) to optimal funds using Sharpe Ratio rankings.'),
    ('Objective 5 - Investor Intelligence: ', 'Perform cohort analysis, SIP continuity tracking, '
     'and demographic segmentation to identify at-risk investors and growth opportunities.'),
], ['Timeline: 7-day intensive capstone project covering data ingestion through final reporting.'])

# ====== SLIDE 3: DATA SOURCES ======
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide3, prs, 'Data Sources & Architecture', [
    ('10 Core Datasets Ingested: ', 'Fund Master (40 schemes), NAV History (64,320 records), '
     'AUM by Fund House (90 records), Monthly SIP Inflows (48 months), Category Inflows (144 records), '
     'Industry Folios (21 snapshots), Investor Transactions (32,778 records), Benchmark Indices (8,050 records).'),
    ('Total Data Volume: ', '127,000+ rows across 17 files including individual fund NAV CSVs '
     'and combined scheme data for 6 major AMCs spanning Jan 2022 to Dec 2025.'),
    ('Star Schema Database: ', 'Fact tables (fact_nav, fact_transactions, fact_performance) + '
     'Dimension tables (dim_fund with 40 schemes, dim_date with 1,608 date entries) loaded '
     'into bluestock_mf.db SQLite database.'),
    ('100% Data Validation: ', 'Automated row-count checks confirmed zero data loss across all 11 tables. '
     'AMFI code cross-validation achieved 100% match rate between Fund Master and NAV History.'),
    ('Tech Stack: ', 'Python 3.14 (Pandas, NumPy, SciPy, Matplotlib, Seaborn, Plotly), '
     'SQLite3, Jupyter Notebooks, Power BI Desktop, Git/GitHub for version control.'),
], ['Data sourced from AMFI (Association of Mutual Funds in India) and supplementary financial data providers.'])

# ====== SLIDE 4: ARCHITECTURE ======
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide4, prs, 'System Architecture: ETL Pipeline', [
    ('EXTRACT (Day 1): ', 'Python script (data_ingestion.py) automatically scans data/raw/ directory, '
     'validates file integrity, detects 17 files with 127,177 total rows, generates a comprehensive '
     'data quality report, and cross-validates AMFI codes across all tables.'),
    ('TRANSFORM (Day 2): ', 'Column standardization to snake_case, datetime parsing, removal of '
     'zero-NAV anomalies, deduplication using keep-last strategy, and dynamic generation of '
     'dim_date table with Year/Quarter/Month/FY flags (1,608 rows).'),
    ('LOAD (Day 2): ', 'Cleaned DataFrames mapped to SQL types and pushed into bluestock_mf.db. '
     '11 tables loaded with automated row-count validation - all achieving 100% match.'),
    ('ANALYZE (Day 3-4-6): ', 'EDA in Jupyter Notebooks using Matplotlib/Seaborn/Plotly. '
     'Performance metrics (CAGR, Sharpe, Sortino, Alpha, Beta) and advanced risk analytics '
     '(VaR, CVaR, HHI, Cohort Analysis) computed programmatically.'),
    ('VISUALIZE (Day 5): ', '4-page interactive Power BI dashboard deployed with KPI cards, '
     'scatter plots, donut charts, heatmaps, dual-axis charts, slicers, and drill-through navigation.'),
    ('DEPLOY (Day 7): ', 'Final report, 12-slide presentation, cleaned GitHub repository with '
     'comprehensive README, and master run_pipeline.py for one-click execution.'),
])

# ====== SLIDE 5: EDA HIGHLIGHTS 1 ======
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide5, prs, 'EDA Highlights: Industry Trends', [
    ('AUM Growth: ', 'Total Industry AUM crossed Rs.81 Lakh Crore by end of 2025, representing '
     'a compound annual growth rate of approximately 18% over the 4-year observation window. '
     'This meteoric rise reflects India\'s deepening capital market participation.'),
    ('Market Concentration: ', 'The Top 10 AMCs (SBI, HDFC, ICICI Prudential, Kotak, Nippon, Axis, '
     'Aditya Birla, UTI, DSP, Mirae Asset) control over 80% of total AUM, indicating '
     'a highly concentrated oligopolistic market structure.'),
    ('SIP Revolution: ', 'Monthly SIP inflows reached an all-time high of Rs.31,002 Crore (Dec 2025). '
     'Active SIP accounts surpassed 26.12 Crore with consistent 15%+ Year-over-Year growth, '
     'confirming a structural behavioral shift from lumpsum to disciplined investing.'),
    ('Folio Explosion: ', 'Total mutual fund folios crossed 26.12 Crore, driven by digital-first '
     'platforms and simplified KYC processes enabling Tier-2/3 city participation.'),
    ('Category Trends: ', 'Equity schemes attracted the lion\'s share of net inflows, while Debt '
     'funds experienced intermittent outflows during rate-hike cycles. ELSS (tax-saver) and '
     'Flexi Cap categories showed the most consistent inflow patterns.'),
])

# ====== SLIDE 6: EDA HIGHLIGHTS 2 ======
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide6, prs, 'EDA Highlights: Investor Demographics', [
    ('Age Distribution: ', 'The 26-35 age cohort dominates SIP investments with the highest '
     'average ticket size (Rs.5,000-10,000/month). The 18-25 cohort shows the fastest adoption '
     'rate but with smaller ticket sizes (Rs.500-2,000/month).'),
    ('Geographic Spread: ', 'Maharashtra and Delhi NCR lead by absolute volume, contributing '
     'over 35% of total transactions. However, Tier-2 cities (Pune, Jaipur, Lucknow) and '
     'Tier-3 cities show 25% higher growth rates, indicating deepening financial inclusion.'),
    ('Transaction Split: ', 'SIPs account for 60%+ of transactions by count, while Lumpsum '
     'investments dominate by total value. Redemptions spike during market correction periods, '
     'particularly among the 56+ age group.'),
    ('Gender Analysis: ', 'Male investors significantly outnumber female investors (approximately '
     '70:30 split), representing a massive untapped market opportunity for targeted campaigns.'),
    ('Payment Modes: ', 'UPI and Net Banking dominate transaction payment modes, with digital '
     'payments accounting for over 85% of all transactions, reflecting India\'s fintech revolution.'),
    ('Income Correlation: ', 'Investors with annual income above Rs.10 Lakh show a strong preference '
     'for Mid Cap and Small Cap funds, while lower-income brackets favor Large Cap and Index funds.'),
])

# ====== SLIDE 7: PERFORMANCE METRICS 1 ======
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide7, prs, 'Performance Metrics: Risk-Adjusted Returns', [
    ('CAGR (Compound Annual Growth Rate): ', 'Computed over 1yr, 3yr, and 5yr horizons using '
     'the formula: CAGR = (NAV_end/NAV_start)^(1/n) - 1. Normalizes returns across different '
     'time periods for fair cross-fund comparison. Top performer: Nippon India Small Cap (5yr CAGR: 28.4%).'),
    ('Sharpe Ratio: ', 'Sharpe = (Rp - Rf) / Std(Rp). Risk-free rate set at 6.5% (RBI repo rate proxy). '
     'Annualized using sqrt(252 trading days). A Sharpe > 1.0 indicates excellent risk-adjusted performance. '
     'Top performer: Mirae Asset Large Cap (Sharpe: 1.45).'),
    ('Sortino Ratio: ', 'Sortino = (Rp - Rf) / Downside_Std. Unlike Sharpe, only penalizes downside '
     'volatility (negative return days). Preferred by conservative investors focused on capital preservation. '
     'Sortino > 1.5 indicates strong downside protection.'),
    ('Alpha & Beta: ', 'Calculated via linear regression (scipy.stats.linregress) against NIFTY 100 benchmark. '
     'Positive Alpha = fund manager skill generating excess returns. Beta > 1.0 = fund is more volatile '
     'than the market. Beta < 1.0 = defensive fund with lower market sensitivity.'),
    ('Maximum Drawdown: ', 'Largest peak-to-trough NAV decline during the observation period. '
     'Critical for worst-case scenario assessment. Funds with drawdowns exceeding 25% were flagged '
     'for heightened volatility risk in the final scorecard.'),
])

# ====== SLIDE 8: PERFORMANCE METRICS 2 ======
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide8, prs, 'Advanced Risk Analytics & Modeling', [
    ('Value at Risk (95% VaR): ', 'Calculated as the 5th percentile of the daily return distribution. '
     'Example: VaR of -2.1% means only a 5% probability of losing more than 2.1% in a single day. '
     'Small Cap funds showed the highest VaR values, confirming higher tail risk.'),
    ('Conditional VaR (CVaR / Expected Shortfall): ', 'Mean of all returns falling below the VaR '
     'threshold. Answers: "When losses exceed VaR, how severe are they on average?" '
     'CVaR provides a more conservative and realistic risk measure than VaR alone.'),
    ('Rolling 90-Day Sharpe Ratio: ', 'Computed for 5 representative funds to visualize time-varying '
     'risk-adjusted performance. Revealed Sharpe compression during mid-2022 global rate hike cycle '
     'followed by strong recovery during the 2023-2024 bull market phases.'),
    ('Sector Concentration (HHI): ', 'Herfindahl-Hirschman Index = sum(weight_i^2). HHI > 0.25 = '
     'highly concentrated portfolio with excessive single-sector exposure. Financial Services '
     'dominated most equity portfolios, with IT as the second-heaviest sector weight.'),
    ('Investor Cohort Analysis: ', 'Segmented investors by first transaction year (2022-2025). '
     'Newer cohorts (2024-2025) invest smaller amounts but with higher frequency. Legacy cohorts (2022) '
     'have larger ticket sizes but show higher attrition rates.'),
    ('SIP Continuity Score: ', 'For investors with 6+ SIP transactions, average gap between payments '
     'was computed. Gaps exceeding 35 days flagged as "at-risk" for discontinuity. This powers '
     'automated retention campaigns to prevent investor churn.'),
])

# ====== SLIDE 9: DASHBOARD PAGE 1 & 2 ======
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_screenshot_slide(slide9, prs, 'Dashboard: Industry Overview (Page 1)',
                     'page1.png',
                     'KPI Cards (Total AUM, SIP Inflows, Folios, Schemes) | AUM Trend Line 2022-2025 | Top 10 AMCs Bar Chart')

# ====== SLIDE 10: DASHBOARD PAGE 3 & 4 ======
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_screenshot_slide(slide10, prs, 'Dashboard: Investor Analytics (Page 3)',
                     'page3.png',
                     'Transaction by State | SIP vs Lumpsum vs Redemption Donut | Age Group vs Avg SIP | Monthly Volume Trend')

# ====== SLIDE 11: KEY FINDINGS ======
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide11, prs, 'Key Findings & Strategic Recommendations', [
    ('Finding 1 - SIP Resilience: ', 'SIP inflows remained structurally resilient even during '
     'market corrections (mid-2022 rate hikes), confirming a permanent behavioral shift in '
     'Indian retail investing. Monthly SIP contributions never dropped below Rs.12,000 Crore.'),
    ('Finding 2 - Risk-Adjusted Leaders: ', 'Mirae Asset and Kotak funds consistently rank '
     'highest on risk-adjusted metrics (Sharpe > 1.3). Some large-cap funds with massive AUM '
     'show negative Alpha, suggesting passive index funds may outperform them.'),
    ('Finding 3 - Tier-2 Opportunity: ', 'Tier-2 city investors show 25% higher SIP growth rates '
     'than Tier-1, presenting a massive untapped distribution and marketing opportunity for fintech platforms.'),
    ('Recommendation 1: ', 'Integrate Sharpe-based fund rankings as default sort order on Bluestock\'s '
     'platform to protect retail investors from chasing high-Beta, high-return traps.'),
    ('Recommendation 2: ', 'Deploy automated SIP churn alerts using the 35-day gap detection logic '
     'to trigger email/SMS retention campaigns for at-risk investors.'),
    ('Recommendation 3: ', 'Expand digital marketing channels targeting 26-35 age group in Tier-2 cities '
     'where customer acquisition cost is lower and growth potential is highest.'),
    ('Recommendation 4: ', 'Display HHI-based "Portfolio Concentration" warnings on fund detail pages '
     'to help investors make informed diversification decisions.'),
])

# ====== SLIDE 12: THANK YOU ======
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide12, prs)

# Top accent
bar_top = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(1.5), Inches(5), Inches(0.04))
bar_top.fill.solid()
bar_top.fill.fore_color.rgb = ACCENT_BLUE
bar_top.line.fill.background()

add_title_textbox(slide12, 'Thank You!', Inches(0.5), Inches(1.8), Inches(9), Inches(1),
                  font_size=44, color=WHITE, alignment=PP_ALIGN.CENTER)

bar_mid = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(3.0), Inches(4), Inches(0.04))
bar_mid.fill.solid()
bar_mid.fill.fore_color.rgb = ACCENT_BLUE
bar_mid.line.fill.background()

add_title_textbox(slide12, 'Bluestock Mutual Fund Analytics Capstone',
                  Inches(0.5), Inches(3.4), Inches(9), Inches(0.6),
                  font_size=18, bold=False, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)
add_title_textbox(slide12, 'Submitted By: Damodara P',
                  Inches(0.5), Inches(4.1), Inches(9), Inches(0.5),
                  font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_title_textbox(slide12, 'Data Analyst Intern  |  Bluestock Fintech  |  June 2026',
                  Inches(0.5), Inches(4.6), Inches(9), Inches(0.5),
                  font_size=14, bold=False, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)
add_title_textbox(slide12, 'pdamodaran2000@gmail.com',
                  Inches(0.5), Inches(5.2), Inches(9), Inches(0.5),
                  font_size=14, bold=False, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

bar_bot = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(5.9), Inches(5), Inches(0.04))
bar_bot.fill.solid()
bar_bot.fill.fore_color.rgb = ACCENT_BLUE
bar_bot.line.fill.background()

add_title_textbox(slide12, 'Questions & Feedback Welcome!',
                  Inches(0.5), Inches(6.2), Inches(9), Inches(0.5),
                  font_size=16, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# ====== SAVE ======
prs.save(OUTPUT_PATH)
print(f"Presentation generated successfully: {OUTPUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
